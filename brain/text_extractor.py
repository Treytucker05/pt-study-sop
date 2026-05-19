"""
Text Extractor — Extract text from PDF, DOCX, PPTX, MD, and TXT files.

PDF extraction uses a tiered approach:
  Tier 1: MinerU  (subprocess, Py 3.12 venv) — best quality, OCR, tables, columns
  Tier 2: PyMuPDF4LLM (direct import) — fast, markdown output, OCR via Tesseract
  Tier 3: pdfplumber (always available) — plain text fallback

Returns: { content: str, error: str | None, metadata: dict }
"""

from __future__ import annotations

import hashlib
import importlib.util
import logging
import inspect
import re
import subprocess
import tempfile
from functools import lru_cache
from pathlib import Path
from typing import Optional

from path_utils import resolve_existing_path

logger = logging.getLogger(__name__)

_PDF_PRIVATE_USE_REPLACEMENTS: dict[str, str] = {
    "\ue062": "Th",
    "\ue0bb": "Th",
    "\uf0e0": " -> ",
}

# ---------------------------------------------------------------------------
# Tier availability checks (cached — only run once per process)
# ---------------------------------------------------------------------------

@lru_cache(maxsize=1)
def _check_pymupdf4llm() -> bool:
    """Return True if pymupdf4llm is importable."""
    try:
        import pymupdf4llm  # noqa: F401
        return True
    except ImportError:
        return False


@lru_cache(maxsize=1)
def _check_pdfplumber() -> bool:
    """Return True if pdfplumber is importable."""
    try:
        import pdfplumber  # noqa: F401

        return True
    except ImportError:
        return False


@lru_cache(maxsize=1)
def _check_docling() -> bool:
    """Return True if Docling is importable."""
    try:
        return importlib.util.find_spec("docling.document_converter") is not None
    except Exception:
        return False


@lru_cache(maxsize=1)
def _check_mineru() -> bool:
    """Return True if the MinerU venv exists with a working mineru executable."""
    repo_root = Path(__file__).resolve().parent.parent
    exe = repo_root / ".venv-mineru" / "Scripts" / "mineru.exe"
    if not exe.exists():
        # Also check Unix-style path
        exe = repo_root / ".venv-mineru" / "bin" / "mineru"
    return exe.exists()


def _mineru_exe() -> Path:
    """Return the path to the MinerU executable."""
    repo_root = Path(__file__).resolve().parent.parent
    exe = repo_root / ".venv-mineru" / "Scripts" / "mineru.exe"
    if exe.exists():
        return exe
    return repo_root / ".venv-mineru" / "bin" / "mineru"


# ---------------------------------------------------------------------------
# Tier 1: MinerU (subprocess to isolated Python 3.12 venv)
# ---------------------------------------------------------------------------

def _extract_pdf_mineru(path: Path) -> str:
    """Extract PDF via MinerU subprocess. Returns markdown string."""
    exe = _mineru_exe()
    with tempfile.TemporaryDirectory() as tmpdir:
        source_arg = str(path)
        if not source_arg.startswith("\\\\?\\"):
            source_arg = source_arg.replace("\\", "/")
        # Use forward slashes for MSYS2 safety
        cmd = [
            str(exe),
            source_arg,
            "-o", tmpdir.replace("\\", "/"),
            "-m", "auto",
        ]
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=300,  # 5 minutes
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"MinerU failed (exit {result.returncode}): {result.stderr[:500]}"
            )

        # MinerU outputs markdown files — find and read them
        md_files = list(Path(tmpdir).rglob("*.md"))
        if not md_files:
            raise RuntimeError("MinerU produced no markdown output")

        parts = []
        for md_file in sorted(md_files):
            parts.append(md_file.read_text(encoding="utf-8", errors="replace"))
        return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Tier 2: PyMuPDF4LLM (direct import — works on Python 3.14)
# ---------------------------------------------------------------------------

def _extract_pdf_pymupdf4llm(path: Path) -> str:
    """Extract PDF via pymupdf4llm. Returns markdown string."""
    import pymupdf4llm

    return pymupdf4llm.to_markdown(str(path))


# ---------------------------------------------------------------------------
# Tier 3: pdfplumber (always available fallback)
# ---------------------------------------------------------------------------

def _extract_pdf_pdfplumber(path: Path) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            # Also extract table text
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    if row:
                        cells = [str(c) if c else "" for c in row]
                        text += "\n" + " | ".join(cells)
            pages.append(text)
    return "\n\n".join(pages)


# ---------------------------------------------------------------------------
# Chapter splitter — for huge textbook PDFs that choke the heavy extractor
# ---------------------------------------------------------------------------

def split_pdf_into_chapters(
    path: str,
    *,
    max_pages: int = 60,
    min_chars: int = 200,
) -> list[dict]:
    """Split a PDF into chapter-sized text segments using its embedded TOC.

    Huge textbooks (100-300 MB, 800-1800 pages) reliably hang the heavy
    tiered extractor (MinerU/OCR on the whole book). But fast page-level
    text extraction (PyMuPDF) per chapter does not. This reads the PDF
    outline, derives top-level (chapter) page ranges, sub-splits any
    range longer than ``max_pages``, and returns one entry per segment::

        [{"title", "text", "page_start", "page_end"}, ...]

    Falls back to fixed ``max_pages`` windows when there is no usable
    outline. Returns ``[]`` on any failure so callers can degrade
    gracefully (skip the file) instead of crashing the sync.
    """
    try:
        import fitz  # PyMuPDF
    except Exception as exc:  # pragma: no cover - env without PyMuPDF
        logger.warning("split_pdf_into_chapters: PyMuPDF unavailable: %s", exc)
        return []

    try:
        doc = fitz.open(path)
    except Exception as exc:
        logger.warning("split_pdf_into_chapters: cannot open %s: %s", path, exc)
        return []

    try:
        page_count = doc.page_count
        if page_count <= 0:
            return []

        # 0-based [start, end) segments from level-1 TOC entries.
        try:
            toc = doc.get_toc(simple=True) or []
        except Exception:
            toc = []
        level1 = [
            (int(entry[2]) - 1, str(entry[1]).strip())
            for entry in toc
            if len(entry) >= 3 and entry[0] == 1 and 1 <= int(entry[2]) <= page_count
        ]

        raw_segments: list[tuple[str, int, int]] = []
        if len(level1) >= 2:
            for i, (start, title) in enumerate(level1):
                end = level1[i + 1][0] if i + 1 < len(level1) else page_count
                if end > start:
                    raw_segments.append((title or f"Section {i + 1}", start, end))
        if not raw_segments:
            # No usable outline → fixed page windows.
            for i, start in enumerate(range(0, page_count, max_pages)):
                end = min(start + max_pages, page_count)
                raw_segments.append((f"Pages {start + 1}-{end}", start, end))

        # Sub-split any segment longer than max_pages so no piece is huge.
        segments: list[tuple[str, int, int]] = []
        for title, start, end in raw_segments:
            if end - start <= max_pages:
                segments.append((title, start, end))
                continue
            part = 1
            for w_start in range(start, end, max_pages):
                w_end = min(w_start + max_pages, end)
                segments.append((f"{title} (part {part})", w_start, w_end))
                part += 1

        results: list[dict] = []
        for title, start, end in segments:
            chunks: list[str] = []
            for pno in range(start, end):
                try:
                    chunks.append(doc.load_page(pno).get_text())
                except Exception:
                    continue
            text = "\n".join(chunks).strip()
            if len(text) < min_chars:
                continue  # cover/blank/figure-only segment — skip
            results.append(
                {
                    "title": title,
                    "text": text,
                    "page_start": start + 1,
                    "page_end": end,
                }
            )
        return results
    finally:
        try:
            doc.close()
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Dispatcher — tries tiers in order, falls through gracefully
# ---------------------------------------------------------------------------

def _has_garbled_content(content: str, threshold: float = 0.05) -> bool:
    """Return True if content is garbled.

    Combines three signals — U+FFFD replacements, Docling GLYPH tags, and
    Latin Extended A/B chars (U+0100-U+024F) — because garbled PDFs often
    spread damage across all three.  Uses 5% threshold on the combined count.
    """
    if not content:
        return False
    n = len(content)
    bad = content.count("\ufffd")
    bad += sum(len(m.group()) for m in re.finditer(r"GLYPH<[^>]*>", content))
    bad += sum(1 for c in content if "\u0100" <= c <= "\u024f")
    return bad / n > threshold


def normalize_extracted_text(content: str) -> str:
    """Clean common embedded-font artifacts from extracted course materials."""
    if not content:
        return content

    text = content
    for glyph, replacement in _PDF_PRIVATE_USE_REPLACEMENTS.items():
        text = text.replace(glyph, replacement)

    # Some PDFs extract ligatures as separated fragments: e ff orts, con fi dence.
    text = re.sub(r"(?<=\w)\s+(ffi|ffl|fi|fl|ff)\s+(?=\w)", r"\1", text)
    text = re.sub(r"\bTh\s+(?=[A-Za-z])", "Th", text)
    text = re.sub(
        r"\b([A-Za-z0-9]+)\s+'\s+(s|t|re|ve|ll|d|m)\b",
        r"\1'\2",
        text,
        flags=re.IGNORECASE,
    )
    text = re.sub(r"@\s+(?=[A-Za-z0-9])", "@", text)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text


def _build_pdf_pipeline_options(*, ocr_mode: bool = False) -> "PdfPipelineOptions":
    """Build Docling PDF pipeline options for normal or OCR extraction."""
    from docling.datamodel.pipeline_options import PdfPipelineOptions, TableStructureOptions

    try:
        from docling.datamodel.pipeline_options import OcrAutoOptions
    except ImportError:
        from docling.datamodel.pipeline_options import OcrOptions as OcrAutoOptions

    table_opts = TableStructureOptions(
        do_cell_matching=not ocr_mode,
        mode="accurate",
    )
    opts = PdfPipelineOptions(
        do_table_structure=True,
        table_structure_options=table_opts,
        generate_picture_images=True,
        images_scale=1.5 if not ocr_mode else 1.0,
    )
    if ocr_mode:
        opts.ocr_options = OcrAutoOptions(force_full_page_ocr=True)
        opts.do_ocr = True
        opts.ocr_batch_size = 1
    else:
        # Keep default extraction lightweight; OCR fallback is handled separately.
        try:
            opts.do_ocr = False
        except Exception:
            pass
    return opts


def _export_docling_markdown(doc: object, source_path: Path) -> str:
    """Export Docling document to markdown, saving images to disk."""
    try:
        from docling_core.types.doc.base import ImageRefMode
    except ImportError:
        ImageRefMode = None

    save_as_markdown = getattr(doc, "save_as_markdown", None)
    if callable(save_as_markdown):
        source_hash = hashlib.md5(str(source_path).encode()).hexdigest()[:12]
        image_dir = Path(__file__).parent / "data" / "extracted_images" / source_hash
        image_dir.mkdir(parents=True, exist_ok=True)

        def _replace_image_placeholders(markdown: str) -> str:
            if "<!-- image -->" not in markdown.lower():
                return markdown
            image_files = sorted(
                [
                    p.name
                    for p in image_dir.iterdir()
                    if p.is_file() and p.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}
                ],
                key=str.lower,
            )
            if not image_files:
                return markdown
            idx = {"value": 0}

            def _replace(_match: re.Match[str]) -> str:
                current = idx["value"]
                idx["value"] += 1
                if current >= len(image_files):
                    return ""
                return f"\n![Extracted image {current + 1}]({image_files[current]})\n"

            return re.sub(r"<!--\s*image\s*-->", _replace, markdown, flags=re.IGNORECASE)

        tmp_md: Optional[Path] = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".md", delete=False) as tmp:
                tmp_md = Path(tmp.name)

            try:
                signature = inspect.signature(save_as_markdown)
                params = set(signature.parameters.keys())
            except Exception:
                params = set()

            kwargs: dict[str, object] = {}
            if "artifacts_dir" in params:
                kwargs["artifacts_dir"] = image_dir
            elif "image_dir" in params:
                kwargs["image_dir"] = image_dir
            elif "assets_dir" in params:
                kwargs["assets_dir"] = image_dir
            if "image_mode" in params and ImageRefMode is not None:
                kwargs["image_mode"] = ImageRefMode.REFERENCED

            attempts: list[dict[str, object]] = [kwargs] if kwargs else []
            if kwargs and "image_mode" in kwargs:
                fallback = dict(kwargs)
                fallback.pop("image_mode", None)
                if fallback not in attempts:
                    attempts.append(fallback)
            if {} not in attempts:
                attempts.append({})

            content = ""
            for attempt in attempts:
                try:
                    save_as_markdown(tmp_md, **attempt)
                    content = tmp_md.read_text(encoding="utf-8")
                    if content.strip():
                        content = _replace_image_placeholders(content)
                        break
                except TypeError:
                    continue

            if content.strip():
                return content
        except Exception as exc:
            logger.debug("save_as_markdown failed, falling back: %s", exc)
        finally:
            if tmp_md is not None:
                tmp_md.unlink(missing_ok=True)
            try:
                if image_dir.exists() and not any(image_dir.iterdir()):
                    image_dir.rmdir()
            except Exception:
                pass

    # Fallback: standard export methods
    for attr in ("export_to_markdown", "export_to_text"):
        exporter = getattr(doc, attr, None)
        if callable(exporter):
            text = exporter()
            if text:
                return str(text)
    for attr in ("text", "raw_text"):
        text = getattr(doc, attr, None)
        if text:
            return str(text)
    raise RuntimeError("Docling conversion produced no content")


@lru_cache(maxsize=None)
def _get_docling_converter(kind: str):
    """Return a process-cached Docling DocumentConverter.

    Rebuilding DocumentConverter per file reloads the ~770-weight
    layout/OCR models every single time — the dominant cost in a bulk
    sync (turned a folder sync into an hours-long crawl). Build one per
    config and reuse it for every file:
      - "pdf"     : PDF layout pipeline, OCR off
      - "pdf_ocr" : PDF pipeline with forced OCR
      - "default" : non-PDF (docx/pptx/etc.)
    """
    from docling.document_converter import DocumentConverter

    if kind == "default":
        return DocumentConverter()

    from docling.datamodel.base_models import InputFormat
    from docling.document_converter import PdfFormatOption

    pipeline_opts = _build_pdf_pipeline_options(ocr_mode=(kind == "pdf_ocr"))
    return DocumentConverter(
        format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_opts)}
    )


def _extract_with_docling(path: Path) -> str:
    """Extract text/markdown using Docling (cached converter)."""
    kind = "pdf" if path.suffix.lower() == ".pdf" else "default"
    converter = _get_docling_converter(kind)
    result = converter.convert(str(path))
    doc = getattr(result, "document", result)
    return _export_docling_markdown(doc, path)


def _docling_ocr_convert(pdf_path: str) -> str:
    """Run Docling forced-OCR on a single PDF (cached converter)."""
    converter = _get_docling_converter("pdf_ocr")
    result = converter.convert(pdf_path)
    doc = getattr(result, "document", result)
    return _export_docling_markdown(doc, Path(pdf_path))


def _subprocess_ocr_convert(pdf_path: str, timeout: int = 300) -> str:
    """Run Docling OCR in a separate subprocess for crash isolation.

    Docling's C++ layer can segfault on certain garbled PDFs, killing the
    host process.  By running OCR in a child process, a segfault (exit 139)
    only kills that child — the parent continues with the next chunk.
    """
    import json
    import sys

    script = (
        "import json, sys; "
        "from pathlib import Path; "
        "sys.path.insert(0, str(Path(__file__).parent) if '__file__' in dir() else '.')\n"
        "from docling.datamodel.base_models import InputFormat; "
        "from docling.document_converter import DocumentConverter, PdfFormatOption; "
        "from docling.datamodel.pipeline_options import PdfPipelineOptions, TableStructureOptions; "
        "try:\n"
        "    from docling.datamodel.pipeline_options import OcrAutoOptions\n"
        "except ImportError:\n"
        "    from docling.datamodel.pipeline_options import OcrOptions as OcrAutoOptions\n"
        "table_opts = TableStructureOptions(do_cell_matching=False, mode='accurate'); "
        "opts = PdfPipelineOptions("
        "do_table_structure=True, table_structure_options=table_opts, "
        "generate_picture_images=False, images_scale=1.0); "
        "opts.ocr_options = OcrAutoOptions(force_full_page_ocr=True); "
        "opts.do_ocr = True; opts.ocr_batch_size = 1; "
        "conv = DocumentConverter("
        "format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}); "
        "r = conv.convert(sys.argv[1]); "
        "doc = getattr(r, 'document', r); "
        "text = ''; "
        "[text := str(f()) for a in ('export_to_markdown','export_to_text') "
        "if (f := getattr(doc, a, None)) and callable(f) and not text]; "
        "json.dump({'text': text or ''}, sys.stdout)"
    )
    try:
        result = subprocess.run(
            [sys.executable, "-c", script, pdf_path],
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        if result.returncode != 0:
            logger.warning(
                "OCR subprocess exited %d for %s: %s",
                result.returncode, pdf_path, result.stderr[:500],
            )
            return ""
        payload = json.loads(result.stdout)
        return payload.get("text", "")
    except subprocess.TimeoutExpired:
        logger.warning("OCR subprocess timed out after %ds for %s", timeout, pdf_path)
        return ""
    except (json.JSONDecodeError, Exception) as exc:
        logger.warning("OCR subprocess result error for %s: %s", pdf_path, exc)
        return ""


def _extract_with_docling_ocr(path: Path, chunk_pages: int = 10) -> str:
    """Extract PDF via Docling with forced full-page OCR.

    Splits large PDFs into chunks and runs each chunk's OCR in a
    separate subprocess for crash isolation (segfault-safe).
    """
    import gc
    import os

    import pymupdf

    src = pymupdf.open(str(path))
    try:
        total = len(src)

        if total <= chunk_pages:
            text = _subprocess_ocr_convert(str(path))
            if not text:
                # Fall back to in-process if subprocess fails
                text = _docling_ocr_convert(str(path))
            if not text:
                raise RuntimeError("Docling OCR conversion produced no content")
            return text

        logger.info("Splitting %d-page PDF into %d-page chunks for OCR (subprocess-isolated)", total, chunk_pages)
        parts: list[str] = []
        for start in range(0, total, chunk_pages):
            end = min(start + chunk_pages, total)
            chunk = pymupdf.open()
            tmp_path = ""
            try:
                chunk.insert_pdf(src, from_page=start, to_page=end - 1)
                with tempfile.NamedTemporaryFile(
                    prefix=f"_ocr_chunk_{start}_", suffix=".pdf", delete=False
                ) as tmp_file:
                    tmp_path = tmp_file.name
                chunk.save(tmp_path)
                chunk.close()
                chunk = None
                text = _subprocess_ocr_convert(tmp_path)
                if not text.strip():
                    # Fall back to in-process for this chunk
                    logger.info("Subprocess empty for pages %d-%d, trying in-process", start, end - 1)
                    text = _docling_ocr_convert(tmp_path)
                if text.strip():
                    parts.append(text)
                    logger.info("OCR chunk pages %d-%d: %d chars", start, end - 1, len(text))
                else:
                    logger.warning("OCR chunk pages %d-%d: empty", start, end - 1)
            except Exception as exc:
                logger.warning("OCR chunk pages %d-%d failed: %s", start, end - 1, exc)
            finally:
                if chunk is not None:
                    try:
                        chunk.close()
                    except Exception:
                        pass
                if tmp_path:
                    try:
                        os.unlink(tmp_path)
                    except OSError:
                        pass
                gc.collect()

        if not parts:
            raise RuntimeError("Docling OCR conversion produced no content")
        return "\n\n".join(parts)
    finally:
        src.close()


def _try_docling(path: Path) -> tuple[Optional[str], list[str]]:
    """Try Docling, retrying with forced OCR if content is mostly garbled."""
    if not _check_docling():
        return None, []
    errors: list[str] = []
    try:
        content = _extract_with_docling(path)
        if not content.strip():
            return None, ["docling: empty content"]

        if _has_garbled_content(content) and path.suffix.lower() == ".pdf":
            logger.warning(
                "Docling text is garbled (broken font mapping) — retrying with forced OCR",
            )
            errors.append("docling: garbled content detected, retried with OCR")
            try:
                ocr_content = _extract_with_docling_ocr(path)
                if ocr_content.strip() and not _has_garbled_content(ocr_content):
                    return ocr_content, errors
                errors.append("docling-ocr: still garbled or empty, using original")
            except Exception as ocr_exc:
                errors.append(f"docling-ocr: {ocr_exc}")

        return content, errors
    except Exception as exc:
        return None, [f"docling: {exc}"]


def _extract_pdf(path: Path) -> tuple[str, str]:
    """
    Try MinerU → PyMuPDF4LLM → pdfplumber.

    Returns (content, extraction_method).
    """
    # Tier 1: MinerU
    if _check_mineru():
        try:
            content = _extract_pdf_mineru(path)
            if content.strip():
                logger.info("PDF extracted via MinerU (Tier 1)")
                return content, "mineru"
            logger.warning("MinerU returned empty content, falling through")
        except Exception as e:
            logger.warning("MinerU extraction failed: %s — falling through", e)

    # Tier 2: PyMuPDF4LLM
    if _check_pymupdf4llm():
        try:
            content = _extract_pdf_pymupdf4llm(path)
            if content.strip():
                logger.info("PDF extracted via PyMuPDF4LLM (Tier 2)")
                return content, "pymupdf4llm"
            logger.warning("PyMuPDF4LLM returned empty content, falling through")
        except Exception as e:
            logger.warning("PyMuPDF4LLM extraction failed: %s — falling through", e)

    # Tier 3: pdfplumber (always available)
    content = _extract_pdf_pdfplumber(path)
    logger.info("PDF extracted via pdfplumber (Tier 3)")
    return content, "pdfplumber"


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_text(file_path: str) -> dict:
    """
    Extract text from a file based on its extension.

    Returns dict with keys: content, error, metadata.
    """
    display_path = Path(file_path)
    path = resolve_existing_path(display_path)
    if not path.exists():
        return {"content": "", "error": f"File not found: {file_path}", "metadata": {}}

    ext = display_path.suffix.lower()
    size = path.stat().st_size

    metadata = {
        "file_name": path.name,
        "file_size": size,
        "file_type": ext.lstrip("."),
    }

    try:
        extraction_errors: list[str] = []
        extractor_name = ""
        extractor_source = ""

        if ext == ".pdf":
            docling_content, extraction_errors = _try_docling(path)
            if docling_content is not None:
                content = docling_content
                extractor_name = "docling"
                extractor_source = "docling"
            else:
                content, extractor_name = _extract_pdf(path)
                extractor_source = "fallback"
        elif ext == ".docx":
            docling_content, extraction_errors = _try_docling(path)
            if docling_content is not None:
                content = docling_content
                extractor_name = "docling"
                extractor_source = "docling"
            else:
                content = _extract_docx(path)
                extractor_name = "python-docx"
                extractor_source = "fallback"
        elif ext == ".pptx":
            docling_content, extraction_errors = _try_docling(path)
            if docling_content is not None:
                content = docling_content
                extractor_name = "docling"
                extractor_source = "docling"
            else:
                content = _extract_pptx(path)
                extractor_name = "python-pptx"
                extractor_source = "fallback"
        elif ext in (".md", ".txt", ".text", ".markdown"):
            content = path.read_text(encoding="utf-8", errors="replace")
        else:
            return {
                "content": "",
                "error": f"Unsupported file type: {ext}",
                "metadata": metadata,
            }

        if ext in (".pdf", ".docx", ".pptx"):
            content = normalize_extracted_text(content)

        if extractor_name:
            metadata["extraction_method"] = extractor_name
            metadata["extractor_name"] = extractor_name
            metadata["extractor_source"] = extractor_source
        if extraction_errors:
            metadata["extraction_errors"] = extraction_errors
        metadata["char_count"] = len(content)
        return {"content": content, "error": None, "metadata": metadata}

    except Exception as e:
        return {"content": "", "error": str(e), "metadata": metadata}


def get_pdf_capabilities() -> dict:
    """Return which PDF extraction tiers are available."""
    return {
        "docling": _check_docling(),
        "mineru": _check_mineru(),
        "pymupdf4llm": _check_pymupdf4llm(),
        "pdfplumber": _check_pdfplumber(),
    }


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    import docx

    doc = docx.Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides):
        slide_text = [f"--- Slide {i + 1} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"[Notes: {notes}]")
        parts.append("\n".join(slide_text))

    return "\n\n".join(parts)


def get_file_type(filename: str) -> Optional[str]:
    """Return normalized file type from filename, or None if unsupported."""
    ext = Path(filename).suffix.lower().lstrip(".")
    mapping = {
        "pdf": "pdf",
        "docx": "docx",
        "pptx": "pptx",
        "md": "md",
        "markdown": "md",
        "txt": "txt",
        "text": "txt",
    }
    return mapping.get(ext)


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".txt", ".text", ".markdown"}
